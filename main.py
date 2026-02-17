import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from data.data import data

def main():
    sprint_number = int(sys.argv[1]) if len(sys.argv) > 1 else 1

    prs = Presentation('./template/base.pptx')

    # Add sprint number to the first slide
    slide = prs.slides[0]
    left = 106 / 72  # points to inches
    top = 310 / 72
    width = 40 / 72
    height = 40 / 72
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    p = textbox.text_frame.paragraphs[0]
    p.text = str(sprint_number)
    p.font.name = 'Arial'
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.bold = True

    # Add slides for each data item
    for slide_data in data:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # blank layout

        # # Add logo
        # left = 870 / 72
        # top = 22 / 72
        # width = 70 / 72
        # height = 13 / 72
        # slide.shapes.add_picture('./images/logo.png', Inches(left), Inches(top), Inches(width), Inches(height))

        # Add epic
        left = 25 / 72
        top = 12 / 72
        width = 800 / 72
        height = 40 / 72
        textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        p = textbox.text_frame.paragraphs[0]
        p.text = slide_data['epic']
        p.font.name = 'Play'
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        # Add title
        left = 25 / 72
        top = 60 / 72
        width = 880 / 72
        height = 60 / 72
        textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        p = textbox.text_frame.paragraphs[0]
        p.text = slide_data['title']
        p.font.name = 'Play'
        p.font.size = Pt(32)
        p.font.color.rgb = RGBColor(255, 255, 255)

        # Add items
        left = 30 / 72
        top = 130 / 72
        width = 450 / 72
        height = 400 / 72
        textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        for item in slide_data['items']:
            parts = item.split('\n')
            for i, part in enumerate(parts):
                p = textbox.text_frame.add_paragraph()
                if i == 0:
                    p.text = f"• {part}"
                    p.level = 0
                else:
                    p.text = part
                    p.level = 1
                p.font.name = 'Arial'
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(255, 255, 255)

    # Move second slide to the end
    current = list(prs.slides._sldIdLst)
    new_order = current[:1] + current[2:] + [current[1]]
    prs.slides._sldIdLst[:] = new_order

    # Ensure result directory exists
    os.makedirs('./result', exist_ok=True)

    prs.save(f'./result/sprint_{sprint_number}.pptx')

if __name__ == "__main__":
    main()
